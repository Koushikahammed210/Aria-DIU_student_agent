# tools.py

import os
import re
import json
import requests
from datetime import datetime
from ddgs import DDGS
from config import SEARCH_MAX_RESULTS, UPLOAD_FOLDER


# ── Web Search ────────────────────────────────────────────────────────────────
def web_search(query: str) -> str:
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=SEARCH_MAX_RESULTS))
        if not results:
            return "No results found."
        output = []
        for i, r in enumerate(results, 1):
            output.append(f"{i}. {r['title']}\n   {r['body']}\n   Source: {r['href']}")
        return "\n\n".join(output)
    except Exception as e:
        return f"Search error: {e}"


# ── Wikipedia ─────────────────────────────────────────────────────────────────
def wikipedia_search(query: str) -> str:
    try:
        resp = requests.get(
            "https://en.wikipedia.org/api/rest_v1/page/summary/" + query.replace(" ", "_"),
            timeout=8
        )
        if resp.status_code == 200:
            data = resp.json()
            return f"{data['title']}\n\n{data['extract']}"
        return "No Wikipedia article found."
    except Exception as e:
        return f"Wikipedia error: {e}"


# ── Weather ───────────────────────────────────────────────────────────────────
def get_weather(location: str) -> str:
    try:
        url = f"https://wttr.in/{requests.utils.quote(location)}?format=3"
        resp = requests.get(url, timeout=8, headers={"User-Agent": "curl/7.0"})
        if resp.status_code == 200:
            return resp.text.strip()
        return f"Could not fetch weather for '{location}'."
    except Exception as e:
        return f"Weather error: {e}"


# ── Summarize Webpage ─────────────────────────────────────────────────────────
def summarize_webpage(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(url, headers=headers, timeout=12)
        text = re.sub(r"<[^>]+>", " ", resp.text)
        text = re.sub(r"\s+", " ", text).strip()
        return f"[Content from {url}]\n\n{text[:4000]}"
    except Exception as e:
        return f"Fetch error: {e}"


# ── Fetch Webpage ─────────────────────────────────────────────────────────────
def fetch_webpage(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(url, headers=headers, timeout=10)
        text = re.sub(r"<[^>]+>", " ", resp.text)
        text = re.sub(r"\s+", " ", text).strip()
        return text[:3000]
    except Exception as e:
        return f"Fetch error: {e}"


# ── File I/O ──────────────────────────────────────────────────────────────────
def read_file(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except FileNotFoundError:
        return f"File not found: {filepath}"
    except Exception as e:
        return f"Error reading file: {e}"


def write_file(filepath: str, content: str) -> str:
    try:
        os.makedirs(os.path.dirname(filepath) or ".", exist_ok=True)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote to {filepath}"
    except Exception as e:
        return f"Error writing file: {e}"


# ── Calculator ────────────────────────────────────────────────────────────────
def calculate(expression: str) -> str:
    try:
        expr = expression.strip()
        match = re.match(r"([\d.]+)\s*%\s*of\s*([\d.]+)", expr, re.IGNORECASE)
        if match:
            a, b = float(match.group(1)), float(match.group(2))
            return f"{a}% of {b} = {a / 100 * b}"
        expr = re.sub(r"([\d.]+)%", r"(\1/100)", expr)
        if not re.match(r"^[\d\s\+\-\*\/\.\(\)]+$", expr):
            return "Error: Only basic math expressions allowed."
        result = eval(expr)
        return f"{expression} = {result}"
    except Exception as e:
        return f"Calculation error: {e}"


# ── Date/Time ─────────────────────────────────────────────────────────────────
def get_datetime() -> str:
    return datetime.now().strftime("%A, %B %d, %Y at %I:%M %p")


# ── File System ───────────────────────────────────────────────────────────────
def list_files(directory: str = ".") -> str:
    try:
        items = os.listdir(directory)
        return "\n".join(items) if items else "Directory is empty."
    except Exception as e:
        return f"Error: {e}"


# ── Run Python ────────────────────────────────────────────────────────────────
def run_python(code: str) -> str:
    import subprocess
    try:
        result = subprocess.run(
            ["python", "-c", code],
            capture_output=True, text=True, timeout=10
        )
        output = result.stdout.strip()
        error = result.stderr.strip()
        if error:
            return f"Error:\n{error}"
        return output if output else "Code ran with no output."
    except subprocess.TimeoutExpired:
        return "Error: Code timed out (10s limit)."
    except Exception as e:
        return f"Run error: {e}"


# ── DIU Cover Page Generator ──────────────────────────────────────────────────
def generate_covers(input_str: str) -> str:
    """
    Generate DIU assignment/lab cover pages (official format) and zip them.

    Expected JSON input:
    {
        "student_name": "...",
        "student_id": "...",
        "section": "...",
        "course_title": "...",
        "course_code": "...",
        "doc_type": "Lab Report" or "Assignment",
        "experiments": ["Exp 1 title", "Exp 2 title"]  -- OR --
                       [{"name": "Exp 1", "no": "1"}, ...],
        "department": "CSE" (optional, default CSE),
        "semester": "6TH" (optional),
        "teacher_name": "..." (optional),
        "teacher_designation": "LECTURER" (optional),
        "teacher_dept": "CSE" (optional),
        "date": "DD/MM/YYYY" (optional, defaults to today)
    }
    """
    try:
        data = json.loads(input_str)
    except json.JSONDecodeError:
        match = re.search(r'\{.*\}', input_str, re.DOTALL)
        if match:
            try:
                data = json.loads(match.group(0))
            except json.JSONDecodeError as e:
                return f"Error: Could not parse input JSON — {e}"
        else:
            return "Error: Input must be a valid JSON object with student info and experiments list."

    required = ["student_name", "student_id", "section", "course_title", "course_code", "experiments"]
    missing = [k for k in required if k not in data]
    if missing:
        return f"Error: Missing required fields: {', '.join(missing)}"

    experiments = data.get("experiments", [])
    if not experiments:
        return "Error: 'experiments' list is empty. Provide at least one experiment name."

    try:
        from cover_generator import batch_create_covers
        student_info = {k: v for k, v in data.items() if k != "experiments"}
        zip_path = batch_create_covers(student_info, experiments)
        filename = os.path.basename(zip_path)
        count = len(experiments)
        return (
            f"SUCCESS: Cover pages generated! "
            f"{count} cover page(s) created. "
            f"Download ready: {filename}"
        )
    except Exception as e:
        return f"Cover generation error: {e}"


# ── File Content Extraction ───────────────────────────────────────────────────
def extract_file_content(filepath: str) -> str:
    """
    Extract text content from uploaded files (PDF, TXT, PPTX, DOCX, etc.)
    Returns extracted text for AI processing.
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".txt" or ext == ".md":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:15000]

        elif ext == ".pdf":
            try:
                import PyPDF2
                text = ""
                with open(filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages[:50]:  # Max 50 pages
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                return text[:15000]
            except ImportError:
                try:
                    import pdfplumber
                    text = ""
                    with pdfplumber.open(filepath) as pdf:
                        for page in pdf.pages[:50]:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                    return text[:15000]
                except ImportError:
                    return "PDF extraction requires PyPDF2 or pdfplumber. Install with: pip install PyPDF2"

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                text = ""
                for i, slide in enumerate(prs.slides, 1):
                    text += f"\n--- Slide {i} ---\n"
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text += para.text + "\n"
                return text[:15000]
            except ImportError:
                return "PPTX extraction requires python-pptx. Install with: pip install python-pptx"

        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(filepath)
                text = ""
                for para in doc.paragraphs:
                    text += para.text + "\n"
                return text[:15000]
            except ImportError:
                return "DOCX extraction requires python-docx. Install with: pip install python-docx"

        elif ext == ".csv":
            try:
                import csv
                text = ""
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text += ", ".join(row) + "\n"
                return text[:15000]
            except Exception as e:
                return f"CSV read error: {e}"

        elif ext == ".json":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
                return json.dumps(data, indent=2)[:15000]

        elif ext in (".py", ".js", ".html", ".css", ".java", ".cpp", ".c"):
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()[:15000]

        else:
            # Try reading as text
            try:
                with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()[:15000]
            except Exception:
                return f"Cannot extract text from {ext} files."

    except Exception as e:
        return f"File extraction error: {e}"


# ── YouTube Video Info ────────────────────────────────────────────────────────
def get_youtube_info(url: str) -> str:
    """Extract basic info from a YouTube URL for tutor mode."""
    try:
        # Extract video ID
        video_id = None
        if "youtube.com/watch" in url:
            import urllib.parse
            parsed = urllib.parse.urlparse(url)
            video_id = urllib.parse.parse_qs(parsed.query).get('v', [None])[0]
        elif "youtu.be/" in url:
            video_id = url.split("youtu.be/")[-1].split("?")[0]

        if not video_id:
            return f"YouTube URL detected but could not extract video ID: {url}"

        # Get video page for title
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(f"https://www.youtube.com/watch?v={video_id}", headers=headers, timeout=10)
        title_match = re.search(r'<title>(.*?)</title>', resp.text)
        title = title_match.group(1).replace(" - YouTube", "") if title_match else "Unknown Title"

        # Try to get transcript/captions info
        return f"YouTube Video: {title}\nURL: {url}\nVideo ID: {video_id}\n\nNote: For detailed learning from this video, share the topic and key points you want to understand."
    except Exception as e:
        return f"YouTube info error: {e}"


# ── Tool Registry ─────────────────────────────────────────────────────────────
TOOLS = {
    "web_search": {
        "fn": web_search,
        "description": "Search the web. Input: search query string.",
    },
    "wikipedia_search": {
        "fn": wikipedia_search,
        "description": "Look up a topic on Wikipedia. Input: topic name.",
    },
    "get_weather": {
        "fn": get_weather,
        "description": "Get current weather for a city. Input: city name (e.g. 'Dhaka' or 'London, UK').",
    },
    "summarize_webpage": {
        "fn": summarize_webpage,
        "description": "Fetch and return full text of a webpage for summarization. Input: full URL.",
    },
    "fetch_webpage": {
        "fn": fetch_webpage,
        "description": "Read raw content of a webpage. Input: full URL.",
    },
    "read_file": {
        "fn": read_file,
        "description": "Read a file from disk. Input: full file path.",
    },
    "write_file": {
        "fn": write_file,
        "description": "Write text to a file. Input: filepath|content (pipe-separated).",
    },
    "calculate": {
        "fn": calculate,
        "description": "Evaluate a math expression. Input: expression like '15% of 3500'.",
    },
    "get_datetime": {
        "fn": get_datetime,
        "description": "Get the current local date and time. Input: none",
    },
    "list_files": {
        "fn": list_files,
        "description": "List files in a directory. Input: directory path or '.'",
    },
    "run_python": {
        "fn": run_python,
        "description": "Run a Python code snippet and return output. Input: valid Python code.",
    },
    "generate_covers": {
        "fn": generate_covers,
        "description": (
            "Generate official DIU-format cover pages (Lab Report or Assignment) as a downloadable ZIP. "
            "Input: JSON string with: student_name, student_id, section, course_title, course_code, "
            'doc_type ("Lab Report" or "Assignment"), '
            "experiments (list of strings OR list of {name, no} objects). "
            "Optional: department, semester, teacher_name, teacher_designation, teacher_dept, date (DD/MM/YYYY). "
        ),
    },
    "get_youtube_info": {
        "fn": get_youtube_info,
        "description": "Get information about a YouTube video from its URL. Input: YouTube URL.",
    },
}
