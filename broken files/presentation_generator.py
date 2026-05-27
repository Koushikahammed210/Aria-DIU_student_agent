# presentation_generator.py — PPTX Presentation Generator for Aria
# FIXED: More robust parsing, never hangs

import os
import json
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Color Palettes ─────────────────────────────────────────────────────────────
PALETTES = {
    "professional": {
        "bg": RGBColor(0x1A, 0x1D, 0x27),
        "accent": RGBColor(0x6C, 0x63, 0xFF),
        "accent2": RGBColor(0xA7, 0x8B, 0xFA),
        "text": RGBColor(0xE2, 0xE8, 0xF0),
        "muted": RGBColor(0x88, 0x92, 0xB0),
        "surface": RGBColor(0x22, 0x26, 0x3A),
    },
    "ocean": {
        "bg": RGBColor(0x0D, 0x23, 0x3B),
        "accent": RGBColor(0x00, 0x96, 0xC7),
        "accent2": RGBColor(0x48, 0xCA, 0xE4),
        "text": RGBColor(0xCA, 0xF0, 0xF8),
        "muted": RGBColor(0x90, 0xE0, 0xEF),
        "surface": RGBColor(0x01, 0x36, 0x5A),
    },
    "forest": {
        "bg": RGBColor(0x0D, 0x1B, 0x0E),
        "accent": RGBColor(0x2D, 0x6A, 0x4F),
        "accent2": RGBColor(0x40, 0x91, 0x6E),
        "text": RGBColor(0xE8, 0xF5, 0xE9),
        "muted": RGBColor(0x81, 0xC7, 0x84),
        "surface": RGBColor(0x1B, 0x34, 0x25),
    },
    "sunset": {
        "bg": RGBColor(0x2D, 0x1B, 0x2E),
        "accent": RGBColor(0xE7, 0x6F, 0x51),
        "accent2": RGBColor(0xF4, 0xA2, 0x61),
        "text": RGBColor(0xFD, 0xE8, 0xD0),
        "muted": RGBColor(0xED, 0xC4, 0xA1),
        "surface": RGBColor(0x4A, 0x25, 0x45),
    },
}


def _add_bg_shape(slide, color, prs_width, prs_height):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), prs_width, prs_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide, color, prs_width):
    """Add a thin accent bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), prs_width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_bullet_slide(slide, title, bullets, palette, prs_width, prs_height):
    """Create a bullet-point content slide."""
    _add_bg_shape(slide, palette["bg"], prs_width, prs_height)
    _add_accent_bar(slide, palette["accent"], prs_width)

    # Title
    _add_text_box(
        slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
        title, font_size=28, color=palette["accent2"], bold=True
    )

    # Separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = palette["accent"]
    sep.line.fill.background()

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Add bullet character
        p.text = f"\u2022  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = palette["text"]
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def _add_title_slide(slide, title, subtitle, palette, prs_width, prs_height):
    """Create the title slide."""
    _add_bg_shape(slide, palette["bg"], prs_width, prs_height)

    # Large accent circle decoration
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.5), Inches(-1.5), Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = palette["surface"]
    circle.line.fill.background()

    # Title
    _add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5),
        title, font_size=36, color=palette["accent2"], bold=True
    )

    # Accent bar under title
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(3.7), Inches(3), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]
    bar.line.fill.background()

    # Subtitle
    _add_text_box(
        slide, Inches(0.8), Inches(4.0), Inches(8.4), Inches(1.2),
        subtitle, font_size=18, color=palette["muted"], bold=False
    )

    # Footer
    _add_text_box(
        slide, Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.5),
        "Generated by Aria AI", font_size=10, color=palette["muted"],
        alignment=PP_ALIGN.LEFT
    )


def _add_content_slide(slide, title, content, palette, prs_width, prs_height):
    """Create a content slide with text body."""
    _add_bg_shape(slide, palette["bg"], prs_width, prs_height)
    _add_accent_bar(slide, palette["accent"], prs_width)

    # Title
    _add_text_box(
        slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
        title, font_size=28, color=palette["accent2"], bold=True
    )

    # Separator
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.15), Inches(2.5), Inches(0.04)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = palette["accent"]
    sep.line.fill.background()

    # Content - handle multi-line content
    lines = content.strip().split("\n")
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check if it's a bullet point
        is_bullet = line.startswith(("- ", "* ", "\u2022 ", "\u2192 "))
        if is_bullet:
            line = line[2:]

        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = palette["text"]
        p.font.name = "Calibri"
        p.space_after = Pt(6)


def _add_section_slide(slide, section_title, palette, prs_width, prs_height):
    """Create a section divider slide."""
    _add_bg_shape(slide, palette["surface"], prs_width, prs_height)

    # Center text
    _add_text_box(
        slide, Inches(1), Inches(2.5), Inches(8), Inches(2),
        section_title, font_size=34, color=palette["accent2"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Accent line
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(4.5), Inches(3), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = palette["accent"]
    bar.line.fill.background()


def _add_end_slide(slide, palette, prs_width, prs_height):
    """Create the thank you / end slide."""
    _add_bg_shape(slide, palette["bg"], prs_width, prs_height)

    _add_text_box(
        slide, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
        "Thank You!", font_size=40, color=palette["accent2"], bold=True,
        alignment=PP_ALIGN.CENTER
    )

    _add_text_box(
        slide, Inches(1), Inches(4.2), Inches(8), Inches(1),
        "Generated by Aria AI", font_size=16, color=palette["muted"],
        alignment=PP_ALIGN.CENTER
    )


def generate_presentation(data: dict) -> str:
    """
    Generate a PPTX presentation from structured data.

    Expected data:
    {
        "topic": "Presentation Title",
        "subtitle": "Optional subtitle",
        "palette": "professional",  # optional
        "slides": [
            {"type": "title", "title": "...", "subtitle": "..."},
            {"type": "bullets", "title": "...", "bullets": ["...", "..."]},
            {"type": "content", "title": "...", "content": "..."},
            {"type": "section", "title": "..."},
        ]
    }

    Returns: path to generated .pptx file
    """
    topic = data.get("topic", "Untitled Presentation")
    subtitle = data.get("subtitle", "")
    palette_name = data.get("palette", "professional")
    palette = PALETTES.get(palette_name, PALETTES["professional"])
    slides_data = data.get("slides", [])

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    prs_width = prs.slide_width
    prs_height = prs.slide_height

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # If no slides provided, create a default structure
    if not slides_data:
        slides_data = [
            {"type": "title", "title": topic, "subtitle": subtitle or "A Presentation"},
        ]

    for slide_info in slides_data:
        try:
            slide = prs.slides.add_slide(blank_layout)
            slide_type = slide_info.get("type", "content")

            if slide_type == "title":
                _add_title_slide(
                    slide,
                    slide_info.get("title", topic),
                    slide_info.get("subtitle", subtitle),
                    palette, prs_width, prs_height
                )
            elif slide_type == "bullets":
                _add_bullet_slide(
                    slide,
                    slide_info.get("title", ""),
                    slide_info.get("bullets", []),
                    palette, prs_width, prs_height
                )
            elif slide_type == "content":
                _add_content_slide(
                    slide,
                    slide_info.get("title", ""),
                    slide_info.get("content", ""),
                    palette, prs_width, prs_height
                )
            elif slide_type == "section":
                _add_section_slide(
                    slide,
                    slide_info.get("title", "Section"),
                    palette, prs_width, prs_height
                )
            elif slide_type == "end":
                _add_end_slide(slide, palette, prs_width, prs_height)
        except Exception as e:
            # Skip broken slides, don't crash the whole generation
            print(f"  [PPTX] Skipping broken slide: {e}")

    # Save the file — try multiple locations
    safe_name = re.sub(r'[^\w\s-]', '', topic).strip()[:50] or "Presentation"
    safe_name = safe_name.replace(' ', '_')
    output_path = os.path.join(os.getcwd(), f"Aria_{safe_name}.pptx")
    prs.save(output_path)
    print(f"  [PPTX] Saved to: {output_path}")
    return output_path


def parse_ai_presentation_response(ai_text: str, topic: str) -> dict:
    """
    Parse the AI's response into structured presentation data.
    IMPROVED: More robust parsing that handles various AI output formats.
    Never crashes — always returns valid data.
    """
    try:
        slides = []
        lines = ai_text.strip().split("\n")

        current_slide = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Detect slide markers — multiple formats supported:
            # "## Slide 1", "# Slide 1", "Slide 1:", "---", "## Title", etc.
            is_new_slide = (
                re.match(r'^#{1,3}\s*(?:Slide|SLIDE)\s*\d', line, re.IGNORECASE) or
                re.match(r'^---+\s*$', line) or
                re.match(r'^Slide\s+\d+', line, re.IGNORECASE) or
                (re.match(r'^#{1,2}\s+', line) and len(line) < 80 and
                 not line.lower().startswith("## title:"))
            )

            if is_new_slide:
                if current_slide:
                    slides.append(current_slide)

                # Extract title from the line
                title = re.sub(r'^#{1,3}\s*(?:Slide|SLIDE)\s*\d*[.:]?\s*', '', line, flags=re.IGNORECASE)
                title = re.sub(r'^---+\s*$', '', title).strip()
                title = re.sub(r'^Slide\s+\d*[.:]?\s*', '', title, flags=re.IGNORECASE).strip()
                current_slide = {"type": "content", "title": title or "Slide"}

            elif line.startswith("Title:") or line.startswith("TITLE:"):
                title = line.split(":", 1)[1].strip()
                if current_slide:
                    current_slide["title"] = title
                else:
                    # Title before any slide marker — create a title slide
                    current_slide = {"type": "title", "title": title, "subtitle": ""}

            elif current_slide and (line.startswith("- ") or line.startswith("* ") or
                                     line.startswith("\u2022 ") or line.startswith("\u2192 ")):
                bullet = line[2:].strip()
                if current_slide.get("type") not in ("bullets", "title"):
                    current_slide["type"] = "bullets"
                    current_slide["bullets"] = []
                if "bullets" in current_slide:
                    current_slide["bullets"].append(bullet)

            elif current_slide and not line.startswith(("#", "---", "Title:", "TITLE:")):
                # Regular content text
                if current_slide.get("type") not in ("bullets", "title"):
                    current_slide["type"] = "content"
                if "content" not in current_slide and current_slide.get("type") == "content":
                    current_slide["content"] = ""
                    current_slide["content"] += line + "\n"

        if current_slide:
            slides.append(current_slide)

        # If no slides parsed, create a simple presentation from the raw text
        if not slides:
            # Try to extract any useful content from the AI response
            paragraphs = [p.strip() for p in ai_text.split("\n\n") if p.strip()]
            slides = [
                {"type": "title", "title": topic, "subtitle": "A Presentation by Aria"},
            ]
            for i, para in enumerate(paragraphs[:6]):
                bullets = [l.strip().lstrip("- *\u2022\u2192") for l in para.split("\n") if l.strip()]
                if bullets:
                    slides.append({"type": "bullets", "title": f"Section {i+1}", "bullets": bullets[:5]})

        # Ensure first slide is a title slide
        if slides and slides[0].get("type") != "title":
            slides.insert(0, {"type": "title", "title": topic, "subtitle": "A Presentation by Aria"})

        # Convert any content-only slides with bullet-like text to bullet slides
        for slide in slides:
            if slide.get("type") == "content" and slide.get("content"):
                content = slide["content"]
                lines_list = [l.strip() for l in content.split("\n") if l.strip()]
                if all(l.startswith(("- ", "* ", "\u2022 ", "\u2192 ")) for l in lines_list if l):
                    slide["type"] = "bullets"
                    slide["bullets"] = [l.lstrip("- *\u2022\u2192").strip() for l in lines_list]
                    slide.pop("content", None)

        # Add end slide if not already there
        if not slides or slides[-1].get("type") != "end":
            slides.append({"type": "end"})

        return {
            "topic": topic,
            "subtitle": "Generated by Aria AI",
            "palette": "professional",
            "slides": slides,
        }

    except Exception as e:
        # ABSOLUTE FALLBACK — never crash
        print(f"  [PPTX Parser error: {e} — using minimal fallback]")
        return {
            "topic": topic,
            "subtitle": "Generated by Aria AI",
            "palette": "professional",
            "slides": [
                {"type": "title", "title": topic, "subtitle": "A Presentation"},
                {"type": "bullets", "title": "Overview", "bullets": [
                    f"Introduction to {topic}",
                    "Key concepts and fundamentals",
                    "Applications and examples",
                    "Summary and conclusions"
                ]},
                {"type": "end"},
            ]
        }
